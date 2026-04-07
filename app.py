import streamlit as st
import requests
from datetime import datetime
import io
import re

# ── Word / PPT / PDF ──────────────────────────────────────────────────────────
from docx import Document
from docx.shared import Pt, RGBColor
from pptx import Presentation
from pptx.util import Inches, Pt as PPTpt
from fpdf import FPDF

# ═════════════════════════════════════════════════════════════════════════════
# PAGE CONFIG
# ═════════════════════════════════════════════════════════════════════════════
st.set_page_config(
    page_title="NLP Chatbot Pro",
    page_icon="🤖",
    layout="wide",
    initial_sidebar_state="expanded"
)

# ═════════════════════════════════════════════════════════════════════════════
# CUSTOM CSS
# ═════════════════════════════════════════════════════════════════════════════
st.markdown("""
<style>
    .stApp { background-color: #0f1117; color: #e0e0e0; }
    section[data-testid="stSidebar"] {
        background: linear-gradient(180deg, #1a1d2e 0%, #12151f 100%);
        border-right: 1px solid #2a2d3e;
    }
    .stButton > button {
        background: linear-gradient(135deg, #667eea, #764ba2);
        color: white; border: none; border-radius: 8px;
        font-weight: 600; transition: opacity 0.2s;
    }
    .stButton > button:hover { opacity: 0.85; color: white; }
    .stDownloadButton > button { border-radius: 8px; font-weight: 600; width: 100%; }
    .feature-header {
        font-size: 1.1rem; font-weight: 700; color: #a78bfa;
        margin: 16px 0 8px 0; padding-bottom: 4px;
        border-bottom: 1px solid #2a2d3e;
    }
    .info-box {
        background: #1e2130; border-left: 3px solid #667eea;
        padding: 10px 14px; border-radius: 6px;
        font-size: 0.85rem; color: #a0aec0; margin-bottom: 12px;
    }
</style>
""", unsafe_allow_html=True)

# ═════════════════════════════════════════════════════════════════════════════
# CONSTANTS
# ═════════════════════════════════════════════════════════════════════════════
api_key         = st.secrets["OPENROUTER_API_KEY"]
CHAT_MODEL      = "openai/gpt-4o-mini"
IMAGE_MODEL     = "openai/dall-e-3"
WHISPER_MODEL   = "openai/whisper-large-v3"
BASE_URL        = "https://openrouter.ai/api/v1"

today = datetime.now().strftime("%A, %B %d, %Y")
SYSTEM_PROMPT = (
    f"You are a helpful, friendly, and knowledgeable assistant. "
    f"Today's date is {today}. Always use this date for date/time questions. "
    f"Give clear, well-structured answers."
)

# ═════════════════════════════════════════════════════════════════════════════
# SESSION STATE
# ═════════════════════════════════════════════════════════════════════════════
if "messages" not in st.session_state:
    st.session_state.messages = []

# ═════════════════════════════════════════════════════════════════════════════
# HELPERS
# ═════════════════════════════════════════════════════════════════════════════

def chat_with_gpt(user_message: str) -> str:
    try:
        payload = [{"role": "system", "content": SYSTEM_PROMPT}] + \
                  st.session_state.messages + \
                  [{"role": "user", "content": user_message}]
        r = requests.post(
            f"{BASE_URL}/chat/completions",
            headers={
                "Authorization": f"Bearer {api_key}",
                "HTTP-Referer": "https://your-app.streamlit.app",
                "X-Title": "NLP Chatbot Pro",
                "Content-Type": "application/json"
            },
            json={"model": CHAT_MODEL, "messages": payload},
            timeout=30
        )
        data = r.json()
        if "choices" in data:
            return data["choices"][0]["message"]["content"]
        elif "error" in data:
            return f"❌ API Error: {data['error'].get('message','Unknown')}"
        return f"❌ Unexpected: {data}"
    except requests.exceptions.Timeout:
        return "❌ Request timed out. Please try again."
    except Exception as e:
        return f"❌ Error: {e}"


def transcribe_audio(audio_bytes: bytes) -> str:
    files = {
        "file": ("audio.wav", io.BytesIO(audio_bytes), "audio/wav"),
        "model": (None, WHISPER_MODEL),
    }
    r = requests.post(
        f"{BASE_URL}/audio/transcriptions",
        headers={"Authorization": f"Bearer {api_key}"},
        files=files,
        timeout=60
    )
    return r.json().get("text", "").strip()


def generate_image(prompt: str, size: str) -> str:
    r = requests.post(
        f"{BASE_URL}/images/generations",
        headers={"Authorization": f"Bearer {api_key}", "Content-Type": "application/json"},
        json={"model": IMAGE_MODEL, "prompt": prompt, "n": 1, "size": size},
        timeout=90
    )
    data = r.json()
    if "data" in data and data["data"]:
        return data["data"][0].get("url", "")
    elif "error" in data:
        raise Exception(data["error"].get("message", "Image error"))
    raise Exception(f"Unexpected: {data}")


def strip_md(text: str) -> str:
    text = re.sub(r'\*\*(.*?)\*\*', r'\1', text)
    text = re.sub(r'\*(.*?)\*', r'\1', text)
    text = re.sub(r'#{1,6}\s*', '', text)
    text = re.sub(r'`{1,3}.*?`{1,3}', '', text, flags=re.DOTALL)
    return text.strip()


# ── Document builders ─────────────────────────────────────────────────────────

def build_docx() -> bytes:
    doc = Document()
    h = doc.add_heading("Chat Export", 0)
    h.runs[0].font.color.rgb = RGBColor(0x66, 0x7e, 0xea)
    p = doc.add_paragraph(f"Exported on: {today}")
    p.runs[0].italic = True
    doc.add_paragraph("")

    for msg in st.session_state.messages:
        label = "🧑 You" if msg["role"] == "user" else "🤖 Assistant"
        rp = doc.add_paragraph()
        run = rp.add_run(label)
        run.bold = True
        run.font.size = Pt(12)
        if msg["role"] == "assistant":
            run.font.color.rgb = RGBColor(0x66, 0x7e, 0xea)
        cp = doc.add_paragraph(strip_md(msg["content"]))
        cp.runs[0].font.size = Pt(11)
        doc.add_paragraph("─" * 50)

    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def build_pptx() -> bytes:
    prs = Presentation()
    # Cover slide
    cover = prs.slides.add_slide(prs.slide_layouts[0])
    cover.shapes.title.text = "Chat Export"
    cover.placeholders[1].text = f"Generated on {today}"

    for msg in st.session_state.messages:
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
        label = "🧑 You" if msg["role"] == "user" else "🤖 Assistant"

        tb_role = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
        tf = tb_role.text_frame
        tf.text = label
        tf.paragraphs[0].runs[0].font.size = PPTpt(16)
        tf.paragraphs[0].runs[0].font.bold = True

        content = strip_md(msg["content"])
        if len(content) > 600:
            content = content[:597] + "..."
        tb_content = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(9), Inches(5.5))
        tf2 = tb_content.text_frame
        tf2.word_wrap = True
        tf2.text = content
        tf2.paragraphs[0].runs[0].font.size = PPTpt(13)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def build_pdf() -> bytes:
    pdf = FPDF()
    pdf.set_auto_page_break(auto=True, margin=15)
    pdf.add_page()

    pdf.set_font("Helvetica", "B", 20)
    pdf.set_text_color(102, 126, 234)
    pdf.cell(0, 12, "Chat Export", ln=True, align="C")

    pdf.set_font("Helvetica", "I", 10)
    pdf.set_text_color(150, 150, 150)
    pdf.cell(0, 8, f"Exported on: {today}", ln=True, align="C")
    pdf.ln(6)

    for msg in st.session_state.messages:
        label = "You:" if msg["role"] == "user" else "Assistant:"
        pdf.set_text_color(60, 180, 120) if msg["role"] == "user" else pdf.set_text_color(102, 126, 234)
        pdf.set_font("Helvetica", "B", 12)
        pdf.cell(0, 8, label, ln=True)

        pdf.set_text_color(40, 40, 40)
        pdf.set_font("Helvetica", "", 11)
        clean = strip_md(msg["content"]).encode("latin-1", errors="replace").decode("latin-1")
        pdf.multi_cell(0, 7, clean)
        pdf.ln(3)
        pdf.set_draw_color(200, 200, 200)
        pdf.line(10, pdf.get_y(), 200, pdf.get_y())
        pdf.ln(4)

    return bytes(pdf.output())


# ═════════════════════════════════════════════════════════════════════════════
# SIDEBAR
# ═════════════════════════════════════════════════════════════════════════════
with st.sidebar:
    st.markdown("## 🤖 NLP Chatbot Pro")
    st.caption(f"📅 {today}")
    st.divider()

    # Image Generation
    st.markdown('<div class="feature-header">🖼️ Image Generation</div>', unsafe_allow_html=True)
    st.markdown('<div class="info-box">Describe any image and DALL·E 3 will create it instantly.</div>', unsafe_allow_html=True)
    img_prompt = st.text_area("Describe the image:", placeholder="A futuristic city at sunset...", height=80, label_visibility="collapsed")
    img_size   = st.selectbox("Size", ["1024x1024", "1024x1792", "1792x1024"], label_visibility="collapsed")

    if st.button("✨ Generate Image", use_container_width=True):
        if img_prompt.strip():
            with st.spinner("Creating image with DALL·E 3..."):
                try:
                    url = generate_image(img_prompt.strip(), img_size)
                    st.image(url, use_container_width=True)
                    st.session_state.messages.append({
                        "role": "assistant",
                        "content": f"🖼️ **Generated image for:** *{img_prompt}*\n\n![image]({url})"
                    })
                    st.success("Image generated and saved to chat!")
                except Exception as e:
                    st.error(f"Image error: {e}")
        else:
            st.warning("Please enter an image description.")

    st.divider()

    # Export Chat
    st.markdown('<div class="feature-header">📄 Export Chat</div>', unsafe_allow_html=True)
    if st.session_state.messages:
        ts = datetime.now().strftime("%Y%m%d_%H%M%S")

        st.download_button(
            "⬇️ Word Document (.docx)",
            data=build_docx(),
            file_name=f"chat_{ts}.docx",
            mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            use_container_width=True
        )
        st.download_button(
            "⬇️ PowerPoint (.pptx)",
            data=build_pptx(),
            file_name=f"chat_{ts}.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            use_container_width=True
        )
        st.download_button(
            "⬇️ PDF Document (.pdf)",
            data=build_pdf(),
            file_name=f"chat_{ts}.pdf",
            mime="application/pdf",
            use_container_width=True
        )

        # Plain text
        plain = f"Chat Export — {today}\n{'='*50}\n\n"
        for m in st.session_state.messages:
            r = "You" if m["role"] == "user" else "Assistant"
            plain += f"{r}:\n{strip_md(m['content'])}\n\n{'─'*40}\n\n"
        st.download_button(
            "⬇️ Plain Text (.txt)",
            data=plain.encode("utf-8"),
            file_name=f"chat_{ts}.txt",
            mime="text/plain",
            use_container_width=True
        )
    else:
        st.caption("💡 Start chatting to enable exports.")

    st.divider()
    if st.button("🗑️ Clear Chat", use_container_width=True):
        st.session_state.messages = []
        st.rerun()

# ═════════════════════════════════════════════════════════════════════════════
# MAIN — Chat History
# ═════════════════════════════════════════════════════════════════════════════
st.markdown("## 💬 Chat")

for message in st.session_state.messages:
    with st.chat_message(message["role"]):
        st.markdown(message["content"])

# ═════════════════════════════════════════════════════════════════════════════
# MAIN — Voice Input
# ═════════════════════════════════════════════════════════════════════════════
st.markdown("---")
st.markdown('<div class="feature-header">🎤 Voice Input</div>', unsafe_allow_html=True)
st.markdown('<div class="info-box">🎙️ Click the mic to record → then press <b>Transcribe & Send</b> to get an answer.</div>', unsafe_allow_html=True)

vcol1, vcol2 = st.columns([4, 1])
with vcol1:
    audio_input = st.audio_input("Record:", label_visibility="collapsed")
with vcol2:
    send_voice = st.button("📝 Transcribe\n& Send", use_container_width=True)

if audio_input and send_voice:
    raw = audio_input.read()
    with st.spinner("🎙️ Transcribing..."):
        try:
            transcribed = transcribe_audio(raw)
            if transcribed:
                st.success(f"**Heard:** {transcribed}")
                st.session_state.messages.append({"role": "user", "content": f"🎤 {transcribed}"})
                with st.spinner("🤖 Thinking..."):
                    reply = chat_with_gpt(transcribed)
                st.session_state.messages.append({"role": "assistant", "content": reply})
                st.rerun()
            else:
                st.warning("Couldn't understand audio. Please speak clearly and try again.")
        except Exception as e:
            st.error(f"Transcription failed: {e}")

# ═════════════════════════════════════════════════════════════════════════════
# MAIN — Text Chat Input
# ═════════════════════════════════════════════════════════════════════════════
st.markdown("---")
user_input = st.chat_input("Type your message here...")

if user_input:
    st.session_state.messages.append({"role": "user", "content": user_input})
    with st.chat_message("user"):
        st.markdown(user_input)

    with st.chat_message("assistant"):
        with st.spinner("Thinking..."):
            reply = chat_with_gpt(user_input)
        st.markdown(reply)

    st.session_state.messages.append({"role": "assistant", "content": reply})
